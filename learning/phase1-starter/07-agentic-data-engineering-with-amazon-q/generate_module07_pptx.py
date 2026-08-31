from __future__ import annotations

import os
import zipfile
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(r"C:\Users\gruddarr\RADE")
MODULE = ROOT / "learning" / "phase1-starter" / "07-agentic-data-engineering-with-amazon-q"
OUTPUT = MODULE / "presentation" / "agentic-data-engineering-with-amazon-q-detailed.pptx"

# Palette requested for the Module 07 deck.
NAVY = "0B1F33"
DARK = "14233E"
SLATE = "556377"
OFF_WHITE = "F4F7FA"
WHITE = "FFFFFF"
TEAL = "009B96"
PURPLE = "6F52B8"
BLUE = "2C66DC"
ORANGE = "E07736"
GREEN = "2A9160"
RED = "C84B4B"
BORDER = "D7E0EA"
PALE_TEAL = "E8F7F5"
PALE_PURPLE = "F1ECFA"
PALE_BLUE = "EAF1FF"
PALE_ORANGE = "FFF1E6"
PALE_GREEN = "EAF6EF"
PALE_RED = "FBECEC"
FONT = "Calibri"

W = 13.333
H = 7.5

CORE = [
    "02-1-intro-aws-setup-revisit.txt",
    "03-2-create-aws-account-revisit.txt",
    "04-3-login-to-aws-using-root-user-and-go-to-iam-revisit.txt",
    "05-35-upgrade-your-aws-account-revisit.txt",
    "06-4-create-admin-user-with-console-and-programmatic-access-revisit.txt",
    "07-5-download-and-install-aws-cli-revisit.txt",
    "08-6-create-access-key-for-aws-cli-access-revisit.txt",
    "09-7-configure-aws-cli-on-your-system-revisit.txt",
    "10-8-vimp-set-up-three-aws-budgets-revisit.txt",
    "11-9-needed-outtro-aws-set-up-revisit.txt",
    "12-1-amazon-q-intro-revisit.txt",
    "13-2-dont-skip-the-last-video-revisit.txt",
    "14-3-kiro-installation-and-setup-by-hema-revisit.txt",
    "15-4-amazon-q-developer-quick-demo-revisit.txt",
    "16-5-requirements-of-the-data-engineering-project-revisit.txt",
    "17-1-steps-in-our-data-engineering-pipeline-use-case-revisit.txt",
    "18-2-create-the-bucket-in-data-lake-with-amazon-q-revisit.txt",
    "19-3-ai-driven-development-vibe-coding-revisit.txt",
    "20-4-below-three-important-use-cases-of-amazon-q-kiro-revisit.txt",
    "21-5-is-ai-taking-away-jobs-where-does-this-leave-us-revisit.txt",
]
SUPPORTING = [
    "25-1-below-mandatory-detour-to-course-rade-agentic-data-engineering-with-amazon-q-r.txt",
    "26-2-before-the-hands-on-continues-vimp-resume.txt",
    "27-3-create-the-glue-ingestion-job-resume.txt",
    "28-4-run-glue-crawler-again-to-update-the-table-metadata-with-partitions-resume.txt",
    "29-5-create-crawler-for-zone-table-resume.txt",
    "30-6-below-next-step-is-very-important-resume.txt",
    "49-7-get-the-script-vetted-by-ai-resume.txt",
]

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def rect(slide, x, y, w, h, fill, line=None, radius=True, transparency=0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(line or fill)
    shape.line.width = Pt(0.8 if line else 0.2)
    return shape


def line(slide, x1, y1, x2, y2, color=BORDER, width=1.0, arrow=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width)
    if arrow:
        conn.line.end_arrowhead = True
    return conn


def text(slide, x, y, w, h, value, size=14, color=DARK, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.03,
         font=FONT, all_caps=False, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    r = p.add_run()
    r.text = value.upper() if all_caps else value
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = rgb(color)
    return box


def rich_text(slide, x, y, w, h, paragraphs, size=14, color=SLATE,
              bullet=False, gap=5, font=FONT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    for i, item in enumerate(paragraphs):
        if isinstance(item, tuple):
            value, is_bold, item_color = item
        else:
            value, is_bold, item_color = item, False, color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if bullet else "") + value
        p.font.name = font
        p.font.size = Pt(size)
        p.font.bold = is_bold
        p.font.color.rgb = rgb(item_color)
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
    return box


def pill(slide, x, y, w, value, fill, color=WHITE, size=9.5):
    rect(slide, x, y, w, 0.30, fill, fill, True)
    text(slide, x + 0.06, y + 0.015, w - 0.12, 0.24, value, size, color, True,
         PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 0, all_caps=True)


def standard_slide(number, category, title_value, subtitle, source, badge="TRANSCRIPT-DERIVED",
                   accent=TEAL):
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, W, H, WHITE, WHITE, False)
    rect(slide, 0, 0, W, 0.12, accent, accent, False)
    text(slide, 0.55, 0.30, 8.7, 0.24, category, 10, accent, True, all_caps=True)
    text(slide, 0.55, 0.61, 11.5, 0.47, title_value, 27, DARK, True)
    text(slide, 0.55, 1.10, 11.7, 0.36, subtitle, 13.5, SLATE)
    text(slide, 12.15, 0.30, 0.62, 0.30, f"{number:02d}", 14, DARK, True, PP_ALIGN.RIGHT)
    badge_width = min(2.30, max(1.38, 0.087 * len(badge) + 0.38))
    badge_fill = RED if "WARNING" in badge or "SAFETY" in badge else accent
    pill(slide, 12.72 - badge_width, 1.25, badge_width, badge, badge_fill, WHITE, 8.1)
    line(slide, 0.55, 7.08, 12.78, 7.08, BORDER, 0.8)
    text(slide, 0.55, 7.14, 9.8, 0.20, source, 6.4, SLATE)
    text(slide, 10.35, 7.14, 2.43, 0.20, "MODULE 07 · AGENTIC DATA ENGINEERING", 6.8, SLATE, True, PP_ALIGN.RIGHT)
    return slide


def section_slide(number, section_no, title_value, promise, source, accent):
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, W, H, NAVY, NAVY, False)
    rect(slide, 0, 0, 0.18, H, accent, accent, False)
    text(slide, 0.70, 0.62, 3.0, 0.28, f"SECTION {section_no}", 11, accent, True, all_caps=True)
    text(slide, 0.70, 1.18, 10.8, 0.92, title_value, 33, WHITE, True)
    pill(slide, 0.70, 2.42, 1.64, "BEGINNER PROMISE", accent, WHITE, 9.0)
    text(slide, 0.70, 2.90, 10.9, 1.35, promise, 18, WHITE, False)
    # Coherent visual: three stepping stones leading to an outcome.
    for i, (label, c) in enumerate([("UNDERSTAND", TEAL), ("DO", BLUE), ("VERIFY", GREEN)]):
        x = 0.70 + i * 2.15
        rect(slide, x, 5.00, 1.70, 0.62, c, c, True)
        text(slide, x, 5.18, 1.70, 0.22, label, 11, WHITE, True, PP_ALIGN.CENTER)
        if i < 2:
            text(slide, x + 1.75, 5.12, 0.35, 0.30, "→", 18, WHITE, True, PP_ALIGN.CENTER)
    text(slide, 12.15, 0.38, 0.62, 0.30, f"{number:02d}", 14, WHITE, True, PP_ALIGN.RIGHT)
    line(slide, 0.70, 6.98, 12.78, 6.98, "34506D", 0.8)
    text(slide, 0.70, 7.08, 9.6, 0.20, source, 6.5, "B9C5D2")
    text(slide, 10.35, 7.08, 2.43, 0.20, "MODULE 07 · AGENTIC DATA ENGINEERING", 6.8, "B9C5D2", True, PP_ALIGN.RIGHT)
    return slide


def card(slide, x, y, w, h, heading, body, accent=TEAL, fill=OFF_WHITE,
         body_size=14, heading_size=15.5, numbered=None, footer=None):
    rect(slide, x, y, w, h, fill, BORDER, True)
    rect(slide, x, y, 0.075, h, accent, accent, False)
    if numbered is not None:
        rect(slide, x + 0.23, y + 0.23, 0.42, 0.42, accent, accent, True)
        text(slide, x + 0.23, y + 0.315, 0.42, 0.20, str(numbered), 11, WHITE, True, PP_ALIGN.CENTER)
        head_x = x + 0.77
        head_w = w - 0.98
    else:
        head_x = x + 0.28
        head_w = w - 0.52
    text(slide, head_x, y + 0.21, head_w, 0.34, heading, heading_size, accent, True, all_caps=True)
    if isinstance(body, str):
        body = [body]
    rich_text(slide, x + 0.25, y + 0.72, w - 0.50, h - 0.91, body, body_size, SLATE, bullet=len(body) > 1, gap=6)
    if footer:
        text(slide, x + 0.25, y + h - 0.35, w - 0.50, 0.20, footer, 9.5, accent, True)


def takeaway(slide, value, accent=TEAL, y=6.25, label="TAKEAWAY"):
    rect(slide, 0.55, y, 12.23, 0.58, NAVY, NAVY, True)
    pill(slide, 0.76, y + 0.14, 1.12, label, accent, WHITE, 8.6)
    text(slide, 2.04, y + 0.15, 10.35, 0.30, value, 13.2, WHITE, True)


def three_cards(slide, items, y=1.78, h=3.95):
    gap = 0.24
    w = (12.23 - 2 * gap) / 3
    for i, item in enumerate(items):
        heading, body, accent, fill = item
        card(slide, 0.55 + i * (w + gap), y, w, h, heading, body, accent, fill)


def two_cards(slide, left, right, y=1.78, h=4.25):
    card(slide, 0.55, y, 5.89, h, *left)
    card(slide, 6.68, y, 6.10, h, *right)


def flow_cards(slide, items, y=2.02, h=2.58, accent=TEAL):
    n = len(items)
    gap = 0.32
    w = (12.23 - gap * (n - 1)) / n
    for i, (heading, body, color) in enumerate(items):
        card(slide, 0.55 + i * (w + gap), y, w, h, heading, body, color, OFF_WHITE, 13.5, 14.5, i + 1)
        if i < n - 1:
            text(slide, 0.55 + (i + 1) * w + i * gap + 0.02, y + 1.02, 0.28, 0.38, "→", 18, accent, True, PP_ALIGN.CENTER)


def checkpoint(slide, items, prompt, accent=GREEN):
    y = 1.80
    for i, item in enumerate(items):
        x = 0.70 + (i % 2) * 6.00
        yy = y + (i // 2) * 1.22
        rect(slide, x, yy, 5.52, 0.90, OFF_WHITE, BORDER, True)
        rect(slide, x + 0.20, yy + 0.20, 0.48, 0.48, accent, accent, True)
        text(slide, x + 0.20, yy + 0.305, 0.48, 0.19, "✓", 12, WHITE, True, PP_ALIGN.CENTER)
        text(slide, x + 0.84, yy + 0.19, 4.38, 0.50, item, 14, DARK, True, valign=MSO_ANCHOR.MIDDLE)
    takeaway(slide, prompt, accent, 6.22, "CHECKPOINT")


def architecture_node(slide, x, y, w, h, heading, sub, color):
    rect(slide, x, y, w, h, color, color, True)
    text(slide, x + 0.15, y + 0.25, w - 0.30, 0.34, heading, 16, WHITE, True, PP_ALIGN.CENTER)
    text(slide, x + 0.15, y + 0.72, w - 0.30, h - 0.83, sub, 12.5, WHITE, False, PP_ALIGN.CENTER)


def comparison_rows(slide, left_title, right_title, rows, left_color=SLATE, right_color=TEAL):
    rect(slide, 0.55, 1.80, 5.88, 0.58, left_color, left_color, True)
    rect(slide, 6.68, 1.80, 6.10, 0.58, right_color, right_color, True)
    text(slide, 0.75, 1.95, 5.48, 0.22, left_title, 14, WHITE, True, PP_ALIGN.CENTER)
    text(slide, 6.88, 1.95, 5.70, 0.22, right_title, 14, WHITE, True, PP_ALIGN.CENTER)
    for i, (lval, rval) in enumerate(rows):
        y = 2.52 + i * 0.84
        rect(slide, 0.55, y, 5.88, 0.65, OFF_WHITE, BORDER, True)
        rect(slide, 6.68, y, 6.10, 0.65, PALE_TEAL, BORDER, True)
        text(slide, 0.78, y + 0.15, 5.42, 0.36, lval, 13.2, DARK, False, valign=MSO_ANCHOR.MIDDLE)
        text(slide, 6.91, y + 0.15, 5.64, 0.36, rval, 13.2, DARK, True, valign=MSO_ANCHOR.MIDDLE)


def index_slide(number, title_value, basenames, source, supporting=None):
    slide = standard_slide(number, "SOURCE INDEX", title_value,
                           "Exact filenames used to establish the evidence boundary.", source,
                           "SOURCE INVENTORY", PURPLE)
    if supporting is None:
        y = 1.70
        for i, basename in enumerate(basenames, 1):
            rect(slide, 0.64, y, 12.05, 0.40, OFF_WHITE, BORDER, True)
            pill(slide, 0.78, y + 0.055, 0.55, f"T{i:02d}", PURPLE, WHITE, 8.2)
            text(slide, 1.48, y + 0.075, 10.92, 0.24, basename, 10.2, DARK, False)
            y += 0.48
    else:
        # Two columns keep all 17 exact basenames visible and clearly separate core from context.
        pill(slide, 0.64, 1.66, 1.53, "CORE · T11–T20", PURPLE, WHITE, 8.2)
        pill(slide, 6.82, 1.66, 2.23, "PROJECT LABS CONTEXT", ORANGE, WHITE, 8.2)
        y_left = 2.08
        for i, basename in enumerate(basenames, 11):
            rect(slide, 0.64, y_left, 5.95, 0.39, OFF_WHITE, BORDER, True)
            pill(slide, 0.76, y_left + 0.045, 0.55, f"T{i:02d}", PURPLE, WHITE, 7.6)
            text(slide, 1.42, y_left + 0.070, 5.00, 0.22, basename, 7.9, DARK, False)
            y_left += 0.45
        y_right = 2.08
        for basename in supporting:
            rect(slide, 6.82, y_right, 5.87, 0.51, PALE_ORANGE, BORDER, True)
            text(slide, 7.02, y_right + 0.10, 5.49, 0.30, basename, 7.7, DARK, False)
            y_right += 0.59
        text(slide, 6.84, 6.33, 5.82, 0.38, "Supporting context only · not counted among the 20 Module 07 core transcripts", 9.2, ORANGE, True, PP_ALIGN.CENTER)
    return slide


# 01 — Cover
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, W, H, NAVY, NAVY, False)
rect(slide, 0, 0, W, 0.15, TEAL, TEAL, False)
text(slide, 0.70, 0.55, 4.0, 0.25, "RADE · MODULE 07", 11, TEAL, True, all_caps=True)
text(slide, 0.70, 1.06, 7.15, 1.36, "Agentic Data Engineering\nwith Amazon Q", 34, WHITE, True)
text(slide, 0.70, 2.68, 6.65, 0.76, "Define the requirement. Direct the agent. Approve the action. Verify the AWS outcome.", 17, "DDE7F0")
pill(slide, 0.70, 3.74, 1.37, "ENGINEER", TEAL)
pill(slide, 2.22, 3.74, 1.14, "AGENT", PURPLE)
pill(slide, 3.51, 3.74, 1.02, "AWS", BLUE)
pill(slide, 4.68, 3.74, 0.94, "S3", GREEN)
# Source → Agent → AWS → S3 visual
nodes = [("NYC TLC\nPARQUET", ORANGE), ("AMAZON Q /\nKIRO", PURPLE), ("AWS GLUE\nPYTHON SHELL", BLUE), ("S3 DATA\nLAKE", GREEN)]
for i, (label, c) in enumerate(nodes):
    x = 7.88 + (i % 2) * 2.18
    y = 1.16 + (i // 2) * 2.02
    rect(slide, x, y, 1.80, 1.18, c, c, True)
    text(slide, x + 0.10, y + 0.33, 1.60, 0.48, label, 14, WHITE, True, PP_ALIGN.CENTER)
text(slide, 9.73, 1.55, 0.30, 0.32, "→", 18, WHITE, True, PP_ALIGN.CENTER)
text(slide, 9.73, 3.57, 0.30, 0.32, "→", 18, WHITE, True, PP_ALIGN.CENTER)
text(slide, 8.56, 2.48, 0.35, 0.34, "↓", 18, WHITE, True, PP_ALIGN.CENTER)
text(slide, 10.74, 2.48, 0.35, 0.34, "↓", 18, WHITE, True, PP_ALIGN.CENTER)
pill(slide, 0.70, 5.32, 2.28, "20 CORE TRANSCRIPTS", "34506D", WHITE, 8.8)
text(slide, 0.70, 5.82, 10.7, 0.55, "Detailed learning deck · editable PowerPoint shapes and text · Project Labs context explicitly separated", 12.5, "B9C5D2")
line(slide, 0.70, 6.93, 12.78, 6.93, "34506D", 0.8)
text(slide, 0.70, 7.06, 10.4, 0.20, f"Sources: {CORE[10]}; {CORE[15]}; {CORE[17]}", 6.5, "B9C5D2")
text(slide, 12.15, 7.02, 0.62, 0.25, "01", 13, WHITE, True, PP_ALIGN.RIGHT)

# 02
slide = standard_slide(2, "ORIENTATION", "How to use this deck—and read its evidence labels",
                       "Learn the workflow, distinguish transcript claims from guidance, then practice with bounded AWS actions.",
                       f"Sources: {CORE[0]}; {CORE[10]}; {CORE[17]}", "TRANSCRIPT + GUIDANCE", TEAL)
three_cards(slide, [
    ("TRANSCRIPT-DERIVED", ["Direct lesson content or a faithful synthesis.", "Footer names the relevant transcript basename."], TEAL, PALE_TEAL),
    ("TRANSCRIPT + GUIDANCE", ["Lesson content plus modern operational guardrails.", "Guidance is not presented as a direct quote."], BLUE, PALE_BLUE),
    ("PROJECT LABS CONTEXT", ["Downstream lessons 25–30 and 49 only.", "Not counted as Module 07 core content."], ORANGE, PALE_ORANGE),
], h=3.90)
takeaway(slide, "Read → pause at checkpoints → perform a small action → capture evidence → explain what happened.", TEAL)

# 03
slide = standard_slide(3, "ORIENTATION", "Learning outcomes: direct, build, diagnose, and explain",
                       "By the end, you should be able to guide an agent without surrendering engineering judgment.",
                       f"Sources: {CORE[10]}; {CORE[15]}; {CORE[17]}; {CORE[18]}; {CORE[19]}", "TRANSCRIPT-DERIVED", TEAL)
flow_cards(slide, [
    ("FRAME", "Turn a business request into a bounded data-engineering requirement.", TEAL),
    ("BUILD", "Use Amazon Q/Kiro to propose code, infrastructure, and commands.", PURPLE),
    ("VERIFY", "Inspect runtime evidence and confirm data exists in S3.", GREEN),
    ("EXPLAIN", "Communicate architecture, trade-offs, and a STAR story.", BLUE),
], y=1.92, h=3.02)
takeaway(slide, "The goal is not faster typing; it is a shorter, safer path from requirement to verified outcome.", TEAL)

# 04
slide = standard_slide(4, "EVIDENCE BOUNDARY", "What is Module 07—and what is only downstream context?",
                       "Twenty TXT files form the core. Seven named Project Labs files provide a clearly labeled handoff.",
                       f"Sources: {CORE[0]}; {CORE[19]}; supporting context: {SUPPORTING[0]}", "SOURCE BOUNDARY", PURPLE)
two_cards(slide,
          ("CORE MODULE 07", ["Exactly 20 TXT files in /transcripts.", "AWS lab setup → agentic workflow → career implications.", "Core content drives slides 1–37 and review slides."], TEAL, PALE_TEAL),
          ("DOWNSTREAM PROJECT LABS", ["Only lessons 25–30 and 49 are used.", "They show the prerequisite gate, ingestion continuation, and later AI code vetting.", "They are not reclassified as Module 07."], ORANGE, PALE_ORANGE),
          h=4.10)
takeaway(slide, "Source boundary rule: context can illuminate the handoff, but it cannot expand the core transcript set.", PURPLE)

# 05
section_slide(5, 1, "Create a safe AWS practice environment",
              "You will set up enough access to learn—while treating root, admin, access keys, and cost controls as deliberate risks.",
              f"Section sources: {CORE[0]}; {CORE[4]}; {CORE[6]}; {CORE[8]}", TEAL)

# 06
slide = standard_slide(6, "SAFE AWS ENVIRONMENT", "AWS setup is a lifecycle, not one checkbox",
                       "The transcript sequence moves from account creation to verified CLI access and budget controls.",
                       f"Sources: {CORE[0]}; {CORE[1]}; {CORE[5]}; {CORE[7]}; {CORE[8]}", "TRANSCRIPT + GUIDANCE", TEAL)
flow_cards(slide, [
    ("ACCOUNT", "Create the learning account and confirm sign-in.", TEAL),
    ("IDENTITY", "Bootstrap IAM access for the course lab.", PURPLE),
    ("CLI", "Install, configure, and verify the AWS CLI.", BLUE),
    ("COST", "Create $10/$20/$30 email budget alerts.", ORANGE),
], y=2.00, h=2.85)
takeaway(slide, "Complete the full lifecycle before asking an agent to create AWS resources.", TEAL)

# 07
slide = standard_slide(7, "SAFE AWS ENVIRONMENT", "Root, course-lab admin, and workplace access are different",
                       "Use the broad course setup only in its intended learning boundary; real organizations provision constrained access.",
                       f"Sources: {CORE[0]}; {CORE[2]}; {CORE[4]}", "CURRENT SAFETY GUIDANCE", RED)
comparison_rows(slide, "COURSE TRANSCRIPT FLOW", "SAFER OPERATING BOUNDARY", [
    ("Root signs in to bootstrap IAM", "Root only for account-level tasks; protect with MFA"),
    ("Create an ADMIN user for practice", "Use a dedicated sandbox—not production data or accounts"),
    ("Attach AdministratorAccess", "Prefer least privilege or an approved temporary role"),
    ("Download sign-in / credential details", "Store secrets securely; rotate or remove when finished"),
])
takeaway(slide, "Broad admin is a course-lab shortcut, not a production IAM pattern.", RED)

# 08
slide = standard_slide(8, "SAFE AWS ENVIRONMENT", "Access keys are credentials—not prompt material",
                       "The secret is shown once in the transcript flow; handle it as a high-impact capability.",
                       f"Source: {CORE[6]}", "CURRENT SAFETY GUIDANCE", RED)
three_cards(slide, [
    ("CREATE DELIBERATELY", ["Create only for the intended CLI identity.", "Record owner, purpose, and expiry."], ORANGE, PALE_ORANGE),
    ("KEEP OUT OF CHAT", ["Never paste access keys, tokens, customer data, or passwords into prompts.", "Use placeholders in examples."], RED, PALE_RED),
    ("PREFER SHORT-LIVED", ["Use approved roles/SSO where available.", "Disable and remove unused long-lived keys."], GREEN, PALE_GREEN),
], h=3.95)
takeaway(slide, "If a command needs a secret, supply it through approved credential mechanisms—not through the agent conversation.", RED)

# 09
slide = standard_slide(9, "SAFE AWS ENVIRONMENT", "Configure the CLI, then prove which identity you are using",
                       "The course uses aws configure and us-east-1; add an identity check before resource creation.",
                       f"Sources: {CORE[5]}; {CORE[7]}", "TRANSCRIPT + GUIDANCE", BLUE)
card(slide, 0.55, 1.78, 7.58, 4.30, "TERMINAL SEQUENCE",
     ["aws --version", "aws configure", "Default region: us-east-1", "aws sts get-caller-identity"],
     BLUE, NAVY, 16, 15.5)
# Restyle code body to white by overlaying deliberate code lines.
rect(slide, 0.88, 2.58, 6.86, 2.72, "102A45", "34506D", True)
text(slide, 1.16, 2.92, 6.30, 1.98, "$ aws --version\n$ aws configure\nDefault region name: us-east-1\n$ aws sts get-caller-identity", 16, WHITE, False)
card(slide, 8.39, 1.78, 4.39, 4.30, "VERIFY BEFORE WRITES",
     ["Account ID matches the intended sandbox.", "Principal/role is expected.", "Region matches the lab.", "No secret is visible in terminal history or chat."],
     GREEN, PALE_GREEN, 13.7)
takeaway(slide, "Configuration success is not identity assurance; verify the caller before an agent runs AWS CLI writes.", BLUE)

# 10
slide = standard_slide(10, "SAFE AWS ENVIRONMENT", "Budgets make forgotten resources visible—not impossible",
                       "Create escalating email alerts at $10, $20, and $30, then still clean up resources after practice.",
                       f"Source: {CORE[8]}", "TRANSCRIPT + GUIDANCE", ORANGE)
three_cards(slide, [
    ("$10 · EARLY SIGNAL", ["Confirm whether spend is expected.", "Inspect active services and regions."], GREEN, PALE_GREEN),
    ("$20 · INVESTIGATE", ["Stop idle jobs and development resources.", "Check data transfer and storage growth."], ORANGE, PALE_ORANGE),
    ("$30 · ACT", ["Escalate cleanup immediately.", "Preserve evidence before deleting resources."], RED, PALE_RED),
], h=3.82)
takeaway(slide, "Budget alerts notify; they do not automatically stop charges. Pair alerts with a cleanup checklist.", ORANGE)

# 11
slide = standard_slide(11, "SAFE AWS ENVIRONMENT", "Setup checkpoint: evidence before automation",
                       "Do the setup, write the steps in your own words, and retain only safe proof.",
                       f"Sources: {CORE[0]}; {CORE[9]}", "CHECKPOINT", GREEN)
checkpoint(slide, [
    "Sandbox account and intended identity confirmed",
    "AWS CLI installed; caller and region verified",
    "$10/$20/$30 budget alerts configured",
    "No credentials copied into notes, screenshots, or prompts",
    "A cleanup plan exists for every created resource",
    "You can explain the setup without reading commands",
], "STOP if identity, region, budget alerts, or credential handling is uncertain.", GREEN)

# 12
section_slide(12, 2, "Frame agentic development before you automate",
              "You will learn a simple mental model: the engineer owns the requirement and judgment; the AI proposes and accelerates execution.",
              f"Section sources: {CORE[10]}; {CORE[12]}; {CORE[13]}; {CORE[15]}", PURPLE)

# 13
slide = standard_slide(13, "AGENTIC DEVELOPMENT", "Engineer/pilot + AI/copilot: one shared control loop",
                       "Amazon Q or Kiro can inspect context and propose actions; the human remains accountable for intent and acceptance.",
                       f"Sources: {CORE[10]}; {CORE[12]}; {CORE[13]}; {CORE[15]}", "TRANSCRIPT + GUIDANCE", PURPLE)
two_cards(slide,
          ("ENGINEER · PILOT", ["Define scope and acceptance criteria.", "Choose what context may be shared.", "Review commands, diffs, cost, and permissions.", "Decide whether evidence proves the outcome."], TEAL, PALE_TEAL),
          ("AI · COPILOT", ["Explain unfamiliar concepts and files.", "Propose code, commands, and infrastructure.", "Read returned logs and suggest fixes.", "Draft documentation and interview prompts."], PURPLE, PALE_PURPLE),
          h=4.15)
takeaway(slide, "Delegating implementation does not delegate accountability.", PURPLE)

# 14
slide = standard_slide(14, "AGENTIC DEVELOPMENT", "Foundations matter more when code arrives faster",
                       "The transcripts explicitly retain SQL, Python/PySpark, architecture, and communication as career-critical skills.",
                       f"Sources: {CORE[10]}; {CORE[15]}; {CORE[19]}", "TRANSCRIPT-DERIVED", TEAL)
flow_cards(slide, [
    ("CONCEPTS", "Know data lakes, Glue, S3, partitions, and execution boundaries.", TEAL),
    ("LANGUAGES", "Read and reason about SQL, Python, and PySpark—even when AI drafts them.", BLUE),
    ("LIFECYCLE", "Understand deployment, logs, CI/CD, and operational ownership.", PURPLE),
    ("COMMUNICATION", "Clarify stakeholder intent and explain decisions in plain language.", ORANGE),
], y=1.92, h=3.14)
takeaway(slide, "AI reduces syntax friction; it increases the value of judgment, architecture, and explanation.", TEAL)

# 15
slide = standard_slide(15, "AGENTIC DEVELOPMENT", "Install Kiro and connect with AWS Builder ID",
                       "Choose the correct OS installer, sign in, open the project folder, and orient yourself before agent mode.",
                       f"Source: {CORE[12]}", "TRANSCRIPT-DERIVED", PURPLE)
flow_cards(slide, [
    ("DOWNLOAD", "Select the Kiro installer for your operating system.", PURPLE),
    ("SIGN IN", "Create or use an AWS Builder ID.", BLUE),
    ("OPEN", "Open the supplied project folder; VS Code import is optional.", TEAL),
    ("ORIENT", "Locate Explorer, editor, chat, and project context.", GREEN),
], y=1.98, h=2.88)
takeaway(slide, "Setup is complete when you can open the intended folder and ask a context-aware question—not merely launch the app.", PURPLE)

# 16
slide = standard_slide(16, "AGENTIC DEVELOPMENT", "Agent mode demo: proposal → approval → filesystem evidence",
                       "The quick demo asks the agent to create project folders and verifies the result in Explorer.",
                       f"Source: {CORE[13]}", "TRANSCRIPT + GUIDANCE", PURPLE)
flow_cards(slide, [
    ("PROMPT", "Create infrastructure and src/glue folders in this project.", TEAL),
    ("PROPOSAL", "Agent shows the shell action it intends to run.", PURPLE),
    ("APPROVAL", "Human checks path and effect, then deliberately runs it.", ORANGE),
    ("EVIDENCE", "Explorer shows the expected folder structure.", GREEN),
], y=1.94, h=3.05)
takeaway(slide, "Chat can explain; agent mode can act. Approve actions only after checking target, scope, and reversibility.", PURPLE)

# 17
slide = standard_slide(17, "AGENTIC DEVELOPMENT", "NYC TLC requirement: real Parquet into an AWS data lake",
                       "The project begins with a source, format, destination, and visible success condition.",
                       f"Source: {CORE[14]}", "TRANSCRIPT-DERIVED", ORANGE)
three_cards(slide, [
    ("SOURCE", ["NYC Taxi & Limousine Commission trip records.", "Yellow, green, FHV, and HVFHV datasets are discussed."], ORANGE, PALE_ORANGE),
    ("FORMAT + DESTINATION", ["Source files are Parquet.", "Land the data in an Amazon S3 data-lake bucket."], BLUE, PALE_BLUE),
    ("ACCEPTANCE", ["AWS prerequisites are ready.", "Files can be seen at the intended S3 path after execution."], GREEN, PALE_GREEN),
], h=3.95)
takeaway(slide, "A useful prompt starts from a testable requirement—not from “build me a pipeline.”", ORANGE)

# 18
slide = standard_slide(18, "AGENTIC DEVELOPMENT", "Connect the nouns: source, Parquet, data lake, S3, and Glue",
                       "A mental model prevents the agent’s generated artifacts from becoming unexplained magic.",
                       f"Sources: {CORE[14]}; {CORE[15]}; {CORE[17]}", "TRANSCRIPT-DERIVED", BLUE)
architecture_node(slide, 0.70, 2.20, 2.10, 1.30, "NYC TLC", "Publishes trip files", ORANGE)
architecture_node(slide, 3.35, 2.20, 2.10, 1.30, "PARQUET", "Columnar source format", PURPLE)
architecture_node(slide, 6.00, 2.20, 2.10, 1.30, "AWS GLUE", "Runs Python ingestion", BLUE)
architecture_node(slide, 8.65, 2.20, 2.10, 1.30, "AMAZON S3", "Durable object storage", GREEN)
architecture_node(slide, 11.30, 2.20, 1.33, 1.30, "DATA LAKE", "Logical purpose", TEAL)
for x in [2.86, 5.51, 8.16, 10.81]:
    text(slide, x, 2.64, 0.42, 0.34, "→", 18, DARK, True, PP_ALIGN.CENTER)
rect(slide, 1.26, 4.30, 10.80, 1.08, PALE_BLUE, BLUE, True)
text(slide, 1.55, 4.56, 10.22, 0.54, "Glue is the execution service. S3 is the storage service. “Data lake” describes how the S3 storage is organized and used.", 15, DARK, True, PP_ALIGN.CENTER)
takeaway(slide, "Name each component’s job before approving code that connects them.", BLUE)

# 19
slide = standard_slide(19, "AGENTIC DEVELOPMENT", "Requirement/prompt checkpoint: constrain the first action",
                       "Turn the lesson requirement into a prompt that exposes assumptions and requests a reviewable plan.",
                       f"Sources: {CORE[14]}; {CORE[15]}; {CORE[16]}", "CHECKPOINT", GREEN)
card(slide, 0.55, 1.76, 7.55, 4.36, "BOUNDED PROMPT",
     ["Goal: land NYC TLC Parquet in an S3 data-lake bucket.", "Region: us-east-1.", "First action: propose a globally unique bucket name and AWS CLI command.", "Before execution: show assumptions, target account, and verification command.", "Do not request or expose credentials."],
     TEAL, PALE_TEAL, 14)
card(slide, 8.36, 1.76, 4.42, 4.36, "HUMAN CHECK",
     ["Correct AWS account?", "Correct region?", "Bucket name valid and unique?", "Command creates only one intended resource?", "Success check defined?"],
     GREEN, PALE_GREEN, 13.5)
takeaway(slide, "Approve the plan only when you can predict the command’s effect and state how you will verify it.", GREEN)

# 20
section_slide(20, 3, "Build, run, diagnose, and verify",
              "You will use the agent to accelerate a real ingestion workflow—then prove the result with runtime and storage evidence.",
              f"Section sources: {CORE[15]}; {CORE[16]}; {CORE[17]}", BLUE)

# 21
slide = standard_slide(21, "BUILD + VERIFY", "The pipeline has three accountable steps",
                       "Create storage, generate ingestion/deployment artifacts, then execute and verify the data outcome.",
                       f"Sources: {CORE[15]}; {CORE[16]}; {CORE[17]}", "TRANSCRIPT-DERIVED", BLUE)
flow_cards(slide, [
    ("CREATE STORAGE", "Create a unique S3 bucket in us-east-1 and confirm it exists.", TEAL),
    ("GENERATE + DEPLOY", "Create Python, CloudFormation YAML, and a deployment shell script.", PURPLE),
    ("RUN + PROVE", "Run Glue, inspect failure/success evidence, and verify objects in S3.", GREEN),
], y=2.05, h=3.00)
takeaway(slide, "Each step needs its own evidence; a later success does not erase an unreviewed earlier action.", BLUE)

# 22
slide = standard_slide(22, "BUILD + VERIFY", "Architecture: CloudFront → Glue Python Shell → S3",
                       "The agent proposes CloudFront as the source path and Glue Python Shell as the execution environment.",
                       f"Source: {CORE[17]}", "TRANSCRIPT-DERIVED", BLUE)
architecture_node(slide, 0.70, 2.16, 2.35, 1.48, "NYC TLC", "Public trip-data publisher", ORANGE)
architecture_node(slide, 3.67, 2.16, 2.35, 1.48, "CLOUDFRONT", "Distributed content endpoint", PURPLE)
architecture_node(slide, 6.64, 2.16, 2.35, 1.48, "GLUE PYTHON SHELL", "Downloads and uploads files", BLUE)
architecture_node(slide, 9.61, 2.16, 2.35, 1.48, "S3 DATA LAKE", "Stores the ingested objects", GREEN)
for x in [3.12, 6.09, 9.06]:
    text(slide, x, 2.72, 0.48, 0.34, "→", 20, DARK, True, PP_ALIGN.CENTER)
rect(slide, 2.10, 4.45, 9.13, 0.92, OFF_WHITE, BORDER, True)
text(slide, 2.34, 4.68, 8.65, 0.40, "Control plane: CloudFormation defines the Glue job; the deployment shell script uploads and deploys artifacts.", 14, DARK, True, PP_ALIGN.CENTER)
takeaway(slide, "Separate data flow from deployment flow so failures can be located precisely.", BLUE)

# 23
slide = standard_slide(23, "BUILD + VERIFY", "Create a unique bucket with an agent-proposed AWS CLI action",
                       "The transcript uses the account number as a uniqueness suffix and verifies the bucket in the console.",
                       f"Source: {CORE[16]}", "TRANSCRIPT + GUIDANCE", TEAL)
card(slide, 0.55, 1.78, 7.25, 4.22, "PROPOSED ACTION",
     ["Resolve the active account context.", "Construct a globally unique bucket name.", "Create the bucket in us-east-1.", "Return the exact verification command and expected result."],
     TEAL, PALE_TEAL, 14.2)
card(slide, 8.05, 1.78, 4.73, 4.22, "APPROVE ONLY IF",
     ["Account and region are intended.", "Name contains no sensitive data.", "No destructive flags are present.", "Tagging/cleanup expectation is known.", "S3 console or CLI proves creation."],
     GREEN, PALE_GREEN, 13.5)
takeaway(slide, "The agent can draft the command; the engineer decides whether this account, region, name, and cost are acceptable.", TEAL)

# 24
slide = standard_slide(24, "BUILD + VERIFY", "Inspect source-access assumptions before writing the downloader",
                       "A public-looking URL is still an interface with availability, naming, and completeness assumptions.",
                       f"Sources: {CORE[14]}; {CORE[17]}", "TRANSCRIPT + GUIDANCE", ORANGE)
three_cards(slide, [
    ("DISCOVERY", ["Which CloudFront URL pattern is being used?", "Which taxi type, year, and month are requested?"], ORANGE, PALE_ORANGE),
    ("TRANSFER", ["Does the runtime have outbound access?", "How are retries, timeouts, and partial downloads handled?"], BLUE, PALE_BLUE),
    ("VALIDATION", ["Expected file extension and naming?", "What object count/size proves completeness?"], GREEN, PALE_GREEN),
], h=3.98)
takeaway(slide, "A downloader is only as reliable as its explicit source contract and completeness checks.", ORANGE)

# 25
slide = standard_slide(25, "BUILD + VERIFY", "Why CloudFront—and why Glue Python Shell?",
                       "Use the source endpoint close to users and a lightweight managed Python runtime for straightforward transfer logic.",
                       f"Source: {CORE[17]}; supporting context: {SUPPORTING[2]}", "TRANSCRIPT + GUIDANCE", BLUE)
two_cards(slide,
          ("CLOUDFRONT", ["The transcript identifies a CloudFront URL for NYC TLC data.", "CloudFront serves content through distributed edge locations.", "The downloader reads the source file from that endpoint."], PURPLE, PALE_PURPLE),
          ("GLUE PYTHON SHELL", ["Runs ordinary Python without a Spark transformation requirement.", "Fits download/upload orchestration for this lesson.", "Later Project Labs context explicitly contrasts it with Spark ETL."], BLUE, PALE_BLUE),
          h=4.15)
takeaway(slide, "Choose the smallest execution model that fits the job; do not invoke Spark merely because the service is Glue.", BLUE)

# 26
slide = standard_slide(26, "BUILD + VERIFY", "Generated artifact map: code, infrastructure, and deployment",
                       "The agent creates three artifacts with different responsibilities—review them separately.",
                       f"Source: {CORE[17]}", "TRANSCRIPT-DERIVED", PURPLE)
three_cards(slide, [
    ("PYTHON", ["Connects to the NYC TLC source.", "Downloads Parquet and uploads it to S3.", "Business/data movement logic."], TEAL, PALE_TEAL),
    ("CLOUDFORMATION YAML", ["Defines the Glue Python Shell job and related configuration.", "Infrastructure as code."], PURPLE, PALE_PURPLE),
    ("DEPLOYMENT SHELL", ["Uploads script artifacts and deploys the stack.", "Operator-facing automation."], ORANGE, PALE_ORANGE),
], h=4.00)
takeaway(slide, "One prompt can produce multiple risk surfaces: data logic, IAM/infrastructure, and command execution.", PURPLE)

# 27
slide = standard_slide(27, "BUILD + VERIFY", "Infrastructure as code vs manual console creation",
                       "Both routes can create the Glue job; IaC makes intent repeatable, diffable, and reviewable.",
                       f"Source: {CORE[17]}", "TRANSCRIPT + GUIDANCE", PURPLE)
comparison_rows(slide, "MANUAL CONSOLE", "CLOUDFORMATION + SCRIPT", [
    ("Open Glue → ETL jobs → script editor", "Review a versioned YAML definition"),
    ("Choose Python Shell and paste code", "Reference/upload the Python artifact"),
    ("Configure and save fields by hand", "Deploy the same declared configuration"),
    ("Harder to reproduce every click", "Repeatable—but only after review and validation"),
])
takeaway(slide, "Automation removes clicks, not responsibility; generated IaC must be understood before deployment.", PURPLE)

# 28
slide = standard_slide(28, "BUILD + VERIFY", "Human review gate: inspect before the command can change AWS",
                       "Review requirement fit, permissions, cost, failure behavior, and verification—not only syntax.",
                       f"Sources: {CORE[13]}; {CORE[17]}; supporting context: {SUPPORTING[6]}", "CURRENT SAFETY GUIDANCE", RED)
checkpoint(slide, [
    "Diff matches the requested source, date scope, and S3 destination",
    "IAM permissions are no broader than the job needs",
    "Commands target the intended account/region and are reversible",
    "Retries, idempotency, logging, and partial failures are considered",
    "No secrets or sensitive data appear in code, YAML, shell, or prompts",
    "Verification checks data—not merely stack/job status",
], "If you cannot explain an artifact, ask the agent to explain it before approval.", RED)

# 29
slide = standard_slide(29, "BUILD + VERIFY", "Deploy the Glue job, run it, and capture run evidence",
                       "The lesson deploys the generated artifacts, starts the job, and checks run status in Glue.",
                       f"Source: {CORE[17]}", "TRANSCRIPT + GUIDANCE", BLUE)
flow_cards(slide, [
    ("DEPLOY", "Run the reviewed deployment script/CloudFormation action.", PURPLE),
    ("CONFIRM", "Open Glue and confirm the expected Python Shell job exists.", BLUE),
    ("RUN", "Start the job with the intended configuration.", ORANGE),
    ("CAPTURE", "Record job run ID, status, timestamps, and log location.", GREEN),
], y=1.92, h=3.15)
takeaway(slide, "A green deployment proves infrastructure creation; it does not prove the ingestion result.", BLUE)

# 30
slide = standard_slide(30, "BUILD + VERIFY", "Evidence loop: failure → logs → diagnosis → change → redeploy",
                       "Return concrete runtime evidence to the agent, then review the proposed correction before rerunning.",
                       f"Source: {CORE[17]}", "TRANSCRIPT + GUIDANCE", ORANGE)
steps = [
    ("FAILURE", "Glue run reports failed", RED),
    ("LOGS", "Capture exact error + run ID", ORANGE),
    ("DIAGNOSIS", "Agent explains likely cause", PURPLE),
    ("CHANGE", "Human reviews focused diff", BLUE),
    ("REDEPLOY", "Run again and compare evidence", GREEN),
]
flow_cards(slide, steps, y=2.00, h=2.92, accent=ORANGE)
takeaway(slide, "Do not report “it failed.” Report the error, context, last known good step, and expected outcome.", ORANGE)

# 31
slide = standard_slide(31, "BUILD + VERIFY", "Verify the S3 outcome: execution success ≠ outcome success",
                       "The transcript checks the bucket after the rerun and observes new folders/files arriving.",
                       f"Source: {CORE[17]}", "TRANSCRIPT + GUIDANCE", GREEN)
two_cards(slide,
          ("RUNTIME EVIDENCE", ["Glue job status is Succeeded.", "Logs reach the terminal condition without hidden errors.", "Run ID and timestamps match the execution you initiated."], BLUE, PALE_BLUE),
          ("DATA OUTCOME EVIDENCE", ["Expected S3 prefix exists.", "Objects have plausible count, size, and Parquet extension.", "A sample object/metadata check matches the requested period and dataset."], GREEN, PALE_GREEN),
          h=4.10)
takeaway(slide, "The acceptance criterion lives in S3: the intended data exists at the intended path and is plausible.", GREEN)

# 32
slide = standard_slide(32, "BUILD + VERIFY", "Vibe coding is evidence-driven iteration—not blind acceptance",
                       "Conversation accelerates build/check/correct cycles only when each cycle returns observable evidence.",
                       f"Source: {CORE[17]}", "TRANSCRIPT + GUIDANCE", PURPLE)
flow_cards(slide, [
    ("STATE INTENT", "Describe the requirement and constraints.", TEAL),
    ("REVIEW PROPOSAL", "Inspect artifacts and commands before execution.", PURPLE),
    ("OBSERVE", "Collect logs, status, files, and data checks.", ORANGE),
    ("REFINE", "Return evidence; approve a focused correction.", GREEN),
], y=1.92, h=3.10)
takeaway(slide, "The “vibe” never outranks the evidence. The loop ends only when the acceptance criteria are verified.", PURPLE)

# 33
slide = standard_slide(33, "BUILD + VERIFY", "Bounded hands-on build checkpoint",
                       "Practice one reversible ingestion slice and retain a small evidence packet you can explain.",
                       f"Sources: {CORE[16]}; {CORE[17]}", "CHECKPOINT", GREEN)
checkpoint(slide, [
    "One unique sandbox bucket created and verified",
    "Source URL and requested Parquet scope documented",
    "Python, YAML, and shell artifacts reviewed separately",
    "One approved deploy/run sequence captured",
    "At least one log/run-status observation retained",
    "S3 prefix, object count/size, and cleanup status verified",
], "Evidence packet: prompt + reviewed diff + command + run ID/log excerpt + S3 proof + cleanup note.", GREEN)

# 34
section_slide(34, 4, "Extend the workflow—and protect your career signal",
              "You will reuse the agent for documentation, interviews, inherited code, and downstream project work without confusing assistance with proof.",
              f"Section sources: {CORE[18]}; {CORE[19]}; supporting context: {SUPPORTING[0]}; {SUPPORTING[6]}", ORANGE)

# 35
slide = standard_slide(35, "EXTEND THE WORKFLOW", "Three practical uses beyond generating pipeline code",
                       "The lesson demonstrates project documentation, interview preparation, and explanation of inherited code.",
                       f"Source: {CORE[18]}", "TRANSCRIPT + GUIDANCE", TEAL)
three_cards(slide, [
    ("DOCUMENTATION", ["Summarize project structure, components, and data flow.", "Reconcile the draft with actual files and deployed resources."], TEAL, PALE_TEAL),
    ("INTERVIEW PREP", ["Generate experience-tailored questions and key points.", "Answer from evidence you personally understand."], BLUE, PALE_BLUE),
    ("INHERITED CODE", ["Explain unfamiliar YAML, architecture, and dependencies.", "Confirm explanations against code and runtime behavior."], PURPLE, PALE_PURPLE),
], h=3.95)
takeaway(slide, "Generated explanations are navigation aids; the repository and runtime remain the source of truth.", TEAL)

# 36
slide = standard_slide(36, "EXTEND THE WORKFLOW", "Use enterprise AI safely: approved context, constrained actions",
                       "Company adoption is growing, but organizational permission does not remove data, access, or review obligations.",
                       f"Sources: {CORE[18]}; {CORE[19]}", "CURRENT SAFETY GUIDANCE", RED)
three_cards(slide, [
    ("APPROVED TOOLING", ["Use the enterprise-approved account, model, and retention policy.", "Know whether code/context can leave the environment."], BLUE, PALE_BLUE),
    ("MINIMUM CONTEXT", ["Share only what is needed.", "Exclude secrets, customer data, credentials, and restricted code."], RED, PALE_RED),
    ("CONTROLLED ACTIONS", ["Prefer read-only discovery first.", "Require human approval for writes, deployments, and deletions."], GREEN, PALE_GREEN),
], h=3.94)
takeaway(slide, "Ask: Is this context approved? Is this action necessary? Can I inspect and reverse the effect?", RED)

# 37
slide = standard_slide(37, "EXTEND THE WORKFLOW", "Career impact: productivity rises—and so does the judgment bar",
                       "The transcript presents two possible outcomes: smaller teams for the same work or more scope delivered faster.",
                       f"Source: {CORE[19]}", "TRANSCRIPT-DERIVED", ORANGE)
two_cards(slide,
          ("WHAT AI COMPRESSES", ["Routine coding and explanation time.", "Project documentation drafting.", "Initial troubleshooting and interview-question generation."], PURPLE, PALE_PURPLE),
          ("WHAT EMPLOYERS VALUE MORE", ["Conceptual maturity and end-to-end architecture.", "Big-picture reasoning and stakeholder communication.", "Data engineering foundations plus CI/CD/DevOps awareness."], ORANGE, PALE_ORANGE),
          h=4.12)
takeaway(slide, "Compete on verified delivery and clear reasoning—not on how much generated code you can produce.", ORANGE)

# 38
slide = standard_slide(38, "PROJECT LABS CONTEXT", "Handoff: Module 07 is an explicit prerequisite boundary",
                       "Lesson 25 pauses the larger project and requires this agentic course before ingestion continues.",
                       f"Supporting context: {SUPPORTING[0]}", "PROJECT LABS CONTEXT", ORANGE)
flow_cards(slide, [
    ("BEFORE", "One Parquet file manually placed in S3; Glue Catalog/crawler and Athena prove access.", SLATE),
    ("NEW ASK", "Load available 2024/2025 trip data plus the taxi-zone lookup.", ORANGE),
    ("MANDATORY DETOUR", "Complete Module 07 to learn AI-assisted source-to-data-lake development.", PURPLE),
    ("RESUME", "Return to Project Labs with the generated ingestion method understood.", GREEN),
], y=1.90, h=3.18)
takeaway(slide, "This is a downstream handoff, not evidence that the whole misrouted folder belongs to Module 07.", ORANGE)

# 39
slide = standard_slide(39, "PROJECT LABS CONTEXT", "Expansion path: partitioned ingestion, metadata, lookup, vetting",
                       "Lessons 26–30 continue the project; lesson 49 later proves that AI-reviewed code still needs human scope checks.",
                       f"Supporting context: {SUPPORTING[1]}; {SUPPORTING[2]}; {SUPPORTING[6]}", "PROJECT LABS CONTEXT", ORANGE)
flow_cards(slide, [
    ("26–27 · INGEST", "Run Glue Python Shell; land 2024 and available 2025 data by year/month; capture CloudWatch logs.", BLUE),
    ("28 · CATALOG", "Rerun the crawler; expose partitions so Athena can prune scans.", PURPLE),
    ("29–30 · ENRICH", "Create the zone table, enable joins, then gather transformation requirements.", TEAL),
    ("49 · VET", "AI enhancement still reads only 2024; human review catches and corrects the missing scope.", RED),
], y=1.88, h=3.38)
takeaway(slide, "AI vetting is one review input—not proof of completeness, correctness, or requirement coverage.", ORANGE)

# 40
slide = standard_slide(40, "INTERVIEW READINESS", "Turn the project into a credible STAR story",
                       "Anchor every claim in a requirement, a decision, runtime evidence, and a verified data result.",
                       f"Sources: {CORE[17]}; {CORE[18]}; {CORE[19]}; supporting context: {SUPPORTING[2]}", "TRANSCRIPT + GUIDANCE", BLUE)
flow_cards(slide, [
    ("SITUATION", "Analysts needed real NYC TLC Parquet in an AWS data lake.", TEAL),
    ("TASK", "Create a repeatable source-to-S3 ingestion path with clear evidence.", BLUE),
    ("ACTION", "Directed Amazon Q/Kiro; reviewed Python/IaC/shell; diagnosed a failed Glue run from logs.", PURPLE),
    ("RESULT", "Verified expected objects in S3 and documented architecture, controls, and lessons.", GREEN),
], y=1.92, h=3.32)
takeaway(slide, "Use real metrics only: object count, date range, run time, scanned data, or cost—never invented impact.", BLUE)

# 41
slide = standard_slide(41, "MODULE REVIEW", "The full operating model in six moves",
                       "Use this sequence whenever an agent helps change data infrastructure or pipeline code.",
                       f"Sources: {CORE[15]}; {CORE[16]}; {CORE[17]}; {CORE[19]}", "MODULE SYNTHESIS", TEAL)
items = [
    ("1 · REQUIRE", "Define source, destination, scope, constraints, and acceptance."),
    ("2 · ORIENT", "Confirm account, identity, region, budget, and approved context."),
    ("3 · PROPOSE", "Ask for a plan, artifacts, commands, and assumptions."),
    ("4 · REVIEW", "Inspect diffs, IAM, cost, secrets, failure modes, and reversibility."),
    ("5 · EXECUTE", "Approve bounded actions and capture run IDs/logs."),
    ("6 · VERIFY", "Check S3/data outcome, document evidence, and clean up."),
]
for i, (head, body) in enumerate(items):
    x = 0.63 + (i % 3) * 4.15
    y = 1.75 + (i // 3) * 2.02
    card(slide, x, y, 3.82, 1.68, head, body, [TEAL, BLUE, PURPLE, ORANGE, GREEN, TEAL][i], OFF_WHITE, 13.3, 14.2)
takeaway(slide, "Engineer/pilot owns the loop. AI/copilot accelerates the steps inside it.", TEAL)

# 42
slide = standard_slide(42, "EVIDENCE NOTES", "Coverage, limitations, normalization, and source boundary",
                       "This deck is detailed learning material—not a claim that every generated action is production-ready.",
                       f"Sources: all 20 core basenames in slides 43–44; supporting context: {SUPPORTING[0]}; {SUPPORTING[6]}", "SOURCE BOUNDARY", PURPLE)
three_cards(slide, [
    ("COVERAGE", ["All 20 core TXT files inform the deck.", "Project Labs use is limited to named lessons 25–30 and 49."], TEAL, PALE_TEAL),
    ("LIMITATIONS", ["Transcript demonstrations are not reproduced screenshots.", "Exact generated code, errors, and object counts are not available in every transcript."], ORANGE, PALE_ORANGE),
    ("NORMALIZATION", ["“Cairo” → Kiro.", "“Amazing Q” → Amazon Q.", "“wipe coding” → vibe coding."], PURPLE, PALE_PURPLE),
], h=3.92)
takeaway(slide, "Safety additions are visibly labeled as TRANSCRIPT + GUIDANCE or CURRENT SAFETY GUIDANCE.", PURPLE)

# 43
index_slide(43, "Core source index · Part 1", CORE[:10], "Sources: " + "; ".join(CORE[:10]))

# 44
slide = index_slide(44, "Core source index · Part 2 + supporting context", CORE[10:],
                    "Core: " + "; ".join(CORE[10:]) + " | Supporting context: " + "; ".join(SUPPORTING),
                    supporting=SUPPORTING)

# Validate expected construction before writing.
assert len(prs.slides) == 44, len(prs.slides)
OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)

# Reopen and validate the generated package.
check = Presentation(OUTPUT)
assert len(check.slides) == 44
assert check.slide_width == Inches(W)
assert check.slide_height == Inches(H)

violations = []
empty_slides = []
prohibited = {"PICTURE", "CHART", "TABLE"}
prohibited_found = []
for slide_no, s in enumerate(check.slides, 1):
    slide_text = []
    for shp in s.shapes:
        if getattr(shp, "has_text_frame", False):
            val = shp.text.strip()
            if val:
                slide_text.append(val)
        if shp.shape_type is not None and str(shp.shape_type).upper() in prohibited:
            prohibited_found.append((slide_no, str(shp.shape_type)))
        # Explicit object-type checks supported by python-pptx.
        if getattr(shp, "has_chart", False):
            prohibited_found.append((slide_no, "CHART"))
        if getattr(shp, "has_table", False):
            prohibited_found.append((slide_no, "TABLE"))
        if shp.left < 0 or shp.top < 0 or shp.left + shp.width > check.slide_width or shp.top + shp.height > check.slide_height:
            violations.append((slide_no, shp.name, shp.left, shp.top, shp.width, shp.height))
    if not slide_text:
        empty_slides.append(slide_no)

# Detect pictures by relationship/content type and verify the zip package.
with zipfile.ZipFile(OUTPUT, "r") as zf:
    bad_zip_member = zf.testzip()
    names = zf.namelist()
    media_members = [n for n in names if n.startswith("ppt/media/")]
    chart_members = [n for n in names if n.startswith("ppt/charts/")]
    embedded_members = [n for n in names if n.startswith("ppt/embeddings/")]

assert not empty_slides, empty_slides
assert not prohibited_found, prohibited_found
assert not violations, violations
assert bad_zip_member is None, bad_zip_member
assert not media_members, media_members
assert not chart_members, chart_members
assert not embedded_members, embedded_members

print(f"OUTPUT={OUTPUT}")
print(f"BYTES={OUTPUT.stat().st_size}")
print(f"SLIDES={len(check.slides)}")
print(f"DIMENSIONS_IN={check.slide_width / 914400:.3f}x{check.slide_height / 914400:.3f}")
print("NONEMPTY_TEXT=PASS")
print("PICTURES_CHARTS_TABLES=PASS")
print("BOUNDS=PASS")
print("ZIP=PASS")
